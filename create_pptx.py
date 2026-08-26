import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re
import codecs

def create_presentation(md_path, pptx_path):
    prs = Presentation()
    
    with codecs.open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
        
    slides_data = content.split('SLIDE ')
    
    for slide_data in slides_data:
        if not slide_data.strip():
            continue
            
        lines = slide_data.strip().split('\n')
        slide_num = lines[0].strip()
        
        # Skip if it's just the intro text before SLIDE 01
        if not slide_num.isdigit() and slide_num != '01':
            continue
            
        title = ""
        body_text = ""
        
        for line in lines[1:]:
            if line.startswith('# '):
                title = line.replace('# ', '').strip()
            elif line.startswith('## '):
                if not title:
                    title = line.replace('## ', '').strip()
                else:
                    body_text += line + "\n"
            else:
                body_text += line + "\n"
                
        # Create a new slide (Title and Content layout)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title if title else f"Slide {slide_num}"
        body_shape.text = body_text.strip()
        
        # Clean up font size to fit more text if needed
        for paragraph in body_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            
    prs.save(pptx_path)

if __name__ == '__main__':
    create_presentation('output/poc_ppt.md', 'output/signal Forge.pptx')
