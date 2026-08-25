import sys
import os
import copy
import asyncio
import glob
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from playwright.async_api import async_playwright

sys.stdout.reconfigure(encoding='utf-8')

def duplicate_slide(prs, source_index):
    source_slide = prs.slides[source_index]
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_el)
    for rel in source_slide.part.rels.values():
        if 'image' in rel.reltype:
            new_slide.part.relate_to(rel.target_part, rel.reltype)
    return new_slide

def clear_all_shapes_except_background(slide):
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame or shape.is_placeholder:
            shapes_to_remove.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            w, h = shape.width.inches, shape.height.inches
            if w > 2.0 and h > 2.0:
                shapes_to_remove.append(shape)
                
    for shape in shapes_to_remove:
        try:
            element = shape.element
            element.getparent().remove(element)
        except Exception:
            pass

async def capture_html_slides():
    html_path = r"file:///d:/workspaces/PMC_POC/output/NextGen_AutoCockpit_Mission_v4.html"
    scratch_dir = r"d:\workspaces\PMC_POC\scratch"
    screenshots = []
    
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page(
            viewport={'width': 1280, 'height': 720}, 
            device_scale_factor=2
        )
        await page.goto(html_path, wait_until='networkidle')
        
        # We know it's a light theme now, but we'll still make background transparent for the PPTX background to show through, OR we can keep the light theme solid.
        # Since the template might be dark or light, wait, the new template is "NextGen AutoCockpit...pptx" which might be dark. 
        # Actually, let's keep the slide-card background as it is (solid white/light) so it looks like a clean card on whatever background.
        await page.add_style_tag(content="""
            body { background: transparent !important; } 
            .container { background: transparent !important; box-shadow: none !important; border: none !important; padding: 0 !important; max-width: none !important;}
        """)
        
        await page.wait_for_timeout(2000)
        cards = await page.query_selector_all('.slide-card')
        
        for i, card in enumerate(cards):
            filename = os.path.join(scratch_dir, f'nextgen_v4_slide_{i:02d}.png')
            await card.screenshot(path=filename, omit_background=True)
            screenshots.append(filename)
            
        await browser.close()
    return screenshots

def build_pptx(screenshots):
    out_dir = r"d:\workspaces\PMC_POC\output"
    
    # Find the template file dynamically to avoid encoding issues
    template_candidates = glob.glob(os.path.join(out_dir, "NextGen AutoCockpit*.pptx"))
    if not template_candidates:
        print("Template file not found.")
        return
        
    template_path = template_candidates[0]
    output_path = os.path.join(out_dir, "NextGen_AutoCockpit_Mission_v4_Presentation.pptx")
    
    print(f"Using template: {template_path}")
    prs = Presentation(template_path)
    
    # We will use slide 1 (index 1) as the master template for content, if exists. 
    # If it has fewer slides, use index 0.
    content_source_idx = 1 if len(prs.slides) > 1 else 0
    
    for idx, img_path in enumerate(screenshots):
        new_slide = duplicate_slide(prs, content_source_idx)
        clear_all_shapes_except_background(new_slide)
        
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        
        margin = Inches(0.4)
        max_w = slide_w - (2 * margin)
        max_h = slide_h - (2 * margin)
        
        pic = new_slide.shapes.add_picture(img_path, 0, 0)
        
        scale_w = max_w / pic.width
        scale_h = max_h / pic.height
        scale = min(scale_w, scale_h)
        
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int((slide_w - pic.width) / 2)
        pic.top = int((slide_h - pic.height) / 2)

    # Remove the original template slides (keep Title slide at index 0)
    xml_slides = prs.slides._sldIdLst
    slides_to_remove = []
    original_count = len(prs.slides) - len(screenshots)
    
    # We want to keep slide 0. So delete from index 1 to original_count - 1
    # But if original_count is 1 (only title slide), we delete nothing.
    for i in range(1, original_count): 
        slides_to_remove.append(xml_slides[i])
        
    for sldId in slides_to_remove:
        xml_slides.remove(sldId)
        
    prs.save(output_path)
    print(f"Successfully generated {output_path} with {len(screenshots)} content slides!")

def main():
    print("Capturing HTML slides...")
    screenshots = asyncio.run(capture_html_slides())
    print(f"Captured {len(screenshots)} screenshots.")
    print("Building PPTX...")
    build_pptx(screenshots)
    
if __name__ == '__main__':
    main()
