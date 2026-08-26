import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette
    C_BG = RGBColor(15, 23, 42)        # #0f172a
    C_SURFACE = RGBColor(30, 41, 59)   # #1e293b
    C_PRIMARY = RGBColor(56, 189, 248)  # #38bdf8
    C_ACCENT = RGBColor(168, 85, 247)  # #a855f7
    C_RED = RGBColor(244, 63, 94)      # #f43f5e
    C_ORANGE = RGBColor(249, 115, 22)  # #f97316
    C_GREEN = RGBColor(34, 197, 94)    # #22c55e
    C_TEXT = RGBColor(248, 250, 252)   # #f8fafc
    C_MUTED = RGBColor(148, 163, 184)  # #94a3b8
    C_BORDER = RGBColor(71, 85, 105)   # #475569

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_BG
        bg.line.color.rgb = C_BG

    def add_header(slide, title_text, category_text="LG SW PM Competition 2025 Final — Stage ① & ② 모범 답안"):
        # Header background banner
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.9))
        banner.fill.solid()
        banner.fill.fore_color.rgb = C_SURFACE
        banner.line.color.rgb = C_ACCENT
        
        tf = banner.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = C_PRIMARY
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(18)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT

    def add_card(slide, left, top, width, height, title, content_list, border_color=C_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = C_SURFACE
        card.line.color.rgb = border_color
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        if title:
            p0 = tf.paragraphs[0]
            p0.text = title
            p0.font.size = Pt(14)
            p0.font.bold = True
            p0.font.color.rgb = C_PRIMARY
            p0.space_after = Pt(8)

        for idx, item in enumerate(content_list):
            p = tf.add_paragraph() if (title or idx > 0) else tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = C_TEXT
            p.space_after = Pt(4)

    # Slide 1: Cover Title
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)
    
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = C_SURFACE
    card1.line.color.rgb = C_ACCENT
    tf1 = card1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "LG SW PM Competition 2025 Final — Case Study 모범 답안"
    p.font.size = Pt(14)
    p.font.color.rgb = C_PRIMARY
    p.font.bold = True
    p.space_after = Pt(14)
    
    p = tf1.add_paragraph()
    p.text = "Stage ① Identify & Stage ② Analyze\n종합 분석 리포트 (v4)"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_TEXT
    p.space_after = Pt(14)
    
    p = tf1.add_paragraph()
    p.text = "시나리오: NovaHome Connect 통합 과제 + 돌발 이슈 01 (내부 갈등: 베타 D-15 소셜로그인 강행)"
    p.font.size = Pt(12)
    p.font.color.rgb = C_MUTED

    # Slide 2: Exec Summary
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "📌 Executive Summary — 프로젝트 위기 진단 요약")
    add_card(slide2, 0.5, 1.5, 6.0, 5.3, "🚨 위기 핵심 상황 (베타 D-15)", [
        "소셜 로그인(REQ-017) 기습 반영으로 인한 기술/빌드 충돌",
        "Nebula SDK 2.4 연동 막바지 자원 부족 및 QA 지원 불능",
        "개인정보 처리 동의 약관 누락 ➔ CRB 승인 거절 법적 리스크",
        "품질 책임자/BE/보안의 전사적 집단 직무 거부 선언",
        "현상은 '일정 지연 및 대립'이나, 근본 원인은 '거버넌스 부재'"
    ], C_RED)
    
    add_card(slide2, 6.8, 1.5, 6.0, 5.3, "🛠️ 5대 핵심 문제 & 3개 실물 도구 일체화", [
        "[문제 1] 거버넌스 붕괴: RACI 표, Stakeholder Matrix, 5 Whys",
        "[문제 2] Scope Creep: RTM 추적표, Sprint 3 Burndown Chart, Mindset",
        "[문제 3] 품질 보증 파이프라인 붕괴: Control Chart, Variance, Fishbone",
        "[문제 4] 법무/보안 Compliance 리스크: Risk 5x5 Heatmap, Quant Risk",
        "[문제 5] 팀 간 R&R 대립: Causal Map, Pareto Chart, Trend Analysis"
    ], C_PRIMARY)

    # Slide 3: Project Info
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "1. 프로젝트 기본 정보 및 게이트 체계 (Identify)")
    add_card(slide3, 0.5, 1.5, 6.0, 5.3, "NovaHome Connect 개요", [
        "프로젝트: 글로벌 가전/스마트홈 통합 플랫폼 구축 (4개월)",
        "론칭 목표일: 2025년 11월 15일 글로벌 동시 론칭",
        "수행 방식: Agile 스크럼 (2주 스프린트) + Milestone Gate (G1~G5)",
        "핵심 파트너: NebulaWorks (인증), AdVantage (광고/마케팅)"
    ])
    add_card(slide3, 6.8, 1.5, 6.0, 5.3, "마일스톤 게이트(Gate) 현황", [
        "G1 Scope Freeze (09/15): 위반 (소셜 로그인 범위 혼선)",
        "G2 Design Freeze (10/01): 위반 (REQ-021 약관 CRB 지연)",
        "G3 Beta (10/20): CRITICAL (D-15 소셜로그인 강행 충돌)",
        "G4 LRR (11/05): CRITICAL (p95 410ms 목표 300ms 초과)",
        "G5 Launch (11/15): High 리스크 미해소 시 론칭 실패 위험"
    ], C_ORANGE)

    # Helper for Problem slides (A & B)
    problems_data = [
        {
            "num": "1",
            "title": "[문제 1] 거버넌스 붕괴 & 권한/스폰서십 비공식 변경",
            "a_items": [
                "Mindset: 의사결정을 공식 체계(CRB)가 아닌 영업/마케팅 압박에 수용함",
                "Domain: Governance Performance Domain & Stakeholder Domain 붕괴",
                "AI Use Case: AI 기반 의사결정 영향도 자동 검증 & RACI 승인 에이전트",
                "Future State: 스폰서-PO-PM 3자 공식 승인 없는 변경 0건 달성"
            ],
            "b_tools": [
                ("Tool 1: RACI Matrix", ["PO가 Accountable(A) 권한 독단 행사", "스폰서/개발팀이 Informed(I)조차 누락된 파괴 확인"]),
                ("Tool 2: Stakeholder Matrix", ["스폰서(전무): Unaware ➔ Leading 갭", "QA/BE: Resistant 상태로 전낙 직무거부"]),
                ("Tool 3: 5 Whys Root Cause", ["Why 1: D-15 소셜로그인 강행 ➔ PO 수용", "Why 5: PMBOK 8판 거버넌스 승인체계 부재"])
            ]
        },
        {
            "num": "2",
            "title": "[문제 2] 요구사항 변경 통제 미비 & Scope Creep",
            "a_items": [
                "Mindset: G1 Scope Freeze 베이스라인 경시 및 의존성 분석 부재",
                "Domain: Planning Performance Domain & Scope Management 붕괴",
                "AI Use Case: AI RTM 추적성 분석기 & 스코프 변동 경고 봇 구축",
                "Future State: 베이스라인 변경 시 자동 영향도 분석 및 스코프크립 0건"
            ],
            "b_tools": [
                ("Tool 1: RTM Traceability", ["REQ-017이 REQ-021약관 및 IAM과 단절", "요구사항 추적성 파괴로 빌드 충돌 발생"]),
                ("Tool 2: Sprint Burndown", ["Day 7에 스코프 +12SP 기습 추가", "번다운 궤적이 위로 꺾이며 완료 위험 입증"]),
                ("Tool 3: Mindset Mapping", ["기존: '요구사항 추가는 흔함, 버티자'", "PMBOK: 'Scope Freeze 준수 & 영향 통제'"])
            ]
        },
        {
            "num": "3",
            "title": "[문제 3] 검증/품질 보증 파이프라인 붕괴",
            "a_items": [
                "Mindset: 캐시 무효화 임시변경이 전체 품질 파이프라인 부채 초래 간과",
                "Domain: Delivery Performance Domain & Quality Performance 붕괴",
                "AI Use Case: AI 자동 테스트 케이스 생성기 & 회귀 성능 예측 모니터링",
                "Future State: CI 빌드 성공률 90% 이상 확보 및 p95 ≤ 300ms 달성"
            ],
            "b_tools": [
                ("Tool 1: Control Chart (관리도)", ["7주 연속 점수 하락 (Rule of Seven 위반)", "7주차 LCL 이탈 (p95 410ms 초과)"]),
                ("Tool 2: Variance Analysis", ["p95 응답속도: 목표 300ms vs 실적 410ms (+110ms)", "CI 성공률: 목표 85% vs 실적 78% (-7%p)"]),
                ("Tool 3: Fishbone Diagram", ["People: QA 40개 TC 전수검증 불능", "Tech: SDK 2.4 빌드충돌 & 스키마 v3 혼선"])
            ]
        },
        {
            "num": "4",
            "title": "[문제 4] 법무/보안 규제 준수(Compliance) 리스크",
            "a_items": [
                "Mindset: 보안/법무 약관을 론칭 직전 형식적 절차로 치부하는 마인드",
                "Domain: Risk Performance Domain & Compliance Domain 붕괴",
                "AI Use Case: AI 컴플라이언스 사전 약관 검증기 & Geo-IP 자동 신청",
                "Future State: CRB 승인 100% 사전 확보 및 법적 리스크 ZERO"
            ],
            "b_tools": [
                ("Tool 1: Risk P-I Heatmap", ["개인정보 동의 약관 누락: P(5) x I(5) = Score 25", "Geo-IP Enforcement 미대응: Score 20"]),
                ("Tool 2: Quant Risk Impact", ["Geo-IP 미신청 시 11/01 유럽 접속 차단", "계약 위반 과징금 + 브랜드 손실 정량 계산"]),
                ("Tool 3: Feedback Loop", ["약관 미비 ➔ CRB 승인 실패 ➔ 론칭 연기", "마케팅 손실 및 전사 신뢰 추락 악순환"])
            ]
        },
        {
            "num": "5",
            "title": "[문제 5] 팀 간 R&R 대립 및 조직 사기 저하",
            "a_items": [
                "Mindset: 리더십을 일방 통제로 인식하고 서번트 리더십 제공 실패",
                "Domain: Team Performance Domain & Culture Domain 붕괴",
                "AI Use Case: AI 다국어 회의록 요약 봇 & 팀 사기 감지 에이전트",
                "Future State: 의사결정 100% 투명 공유 및 팀 협업 만족도 90% 달성"
            ],
            "b_tools": [
                ("Tool 1: Causal Relationship Map", ["거버넌스 붕괴 ➔ Scope Creep ➔ 품질 붕괴", "개발/QA/보안 책임 전가 ➔ 이탈 위험 악순환"]),
                ("Tool 2: Pareto Chart (80/20)", ["불만 원인의 60.8%가 '일방 스코프 변경'", "기술 부채 강요 2개 상위 원인에 집중됨"]),
                ("Tool 3: Trend Analysis", ["인도 QA Raj 회의록 오역(Disabled vs Delay)", "블라인드 게시판 불만 폭발 ➔ 이탈 추세"])
            ]
        }
    ]

    for p in problems_data:
        # Slide A: Analysis
        slide_a = prs.slides.add_slide(blank_layout)
        add_bg(slide_a)
        add_header(slide_a, f"🔍 {p['title']} (분석)")
        add_card(slide_a, 0.5, 1.5, 12.333, 5.3, "PMBOK 8판 4단계 논리 분석 체계", p["a_items"], C_ACCENT)

        # Slide B: 3 Tools Full Package
        slide_b = prs.slides.add_slide(blank_layout)
        add_bg(slide_b)
        add_header(slide_b, f"📊 {p['title']} — 실물 디텍팅 도구 3종 풀패키지")
        
        # 3 Cards side-by-side
        add_card(slide_b, 0.5, 1.5, 3.8, 5.3, p["b_tools"][0][0], p["b_tools"][0][1], C_PRIMARY)
        add_card(slide_b, 4.76, 1.5, 3.8, 5.3, p["b_tools"][1][0], p["b_tools"][1][1], C_PRIMARY)
        add_card(slide_b, 9.03, 1.5, 3.8, 5.3, p["b_tools"][2][0], p["b_tools"][2][1], C_PRIMARY)

    # Final Conclusion Slide
    slide_final = prs.slides.add_slide(blank_layout)
    add_bg(slide_final)
    add_header(slide_final, "🚀 Stage ① & ② 종합 결론 및 Stage ③ Design 이행 로드맵")
    add_card(slide_final, 0.5, 1.5, 12.333, 5.3, "종합 진단 총평 & Next Step", [
        "5개 문제에 매핑된 실물 분석 도구 15종(RACI, RTM, Control Chart, Burndown, Pareto 등)을 통해",
        "증상 뒤에 숨겨진 'PMBOK 8판 거버넌스 붕괴 및 비공식 변경 절차'라는 근본 원인을 완벽히 증명했습니다.",
        "Stage ① Identify: 5대 KPI 갭 식별 완료 (p95 410ms, CI 78%, 돌발 이슈 01 포착)",
        "Stage ② Analyze: 문제별 3개 실물 도구(총 15종)를 내장하여 근본 원인 증명 완료",
        "Stage ③ Design 이행: AI RTM Agent, CRB 자동검증 봇, AI 회의록 시스템 기반 개선안 설계 착수"
    ], C_GREEN)

    prs.save(output_path)
    print(f"Successfully generated PPTX: {output_path}")

if __name__ == "__main__":
    out1 = r"d:\workspaces\PMC_POC\stage1_stage2_analysis_v4.pptx"
    out2 = r"d:\workspaces\PMC_POC\output\stage1_stage2_analysis_v4.pptx"
    create_deck(out1)
    create_deck(out2)
