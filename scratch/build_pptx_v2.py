import sys
import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

sys.stdout.reconfigure(encoding='utf-8')

# Color constants matching Dark Tech Theme
COLOR_BG_DARK = RGBColor(15, 23, 42)      # #0f172a
COLOR_CARD = RGBColor(30, 41, 59)        # #1e293b
COLOR_CARD_ALT = RGBColor(51, 65, 85)    # #334155
COLOR_CYAN = RGBColor(56, 189, 248)      # #38bdf8
COLOR_PURPLE = RGBColor(168, 85, 247)    # #a855f7
COLOR_RED = RGBColor(244, 63, 94)        # #f43f5e
COLOR_ORANGE = RGBColor(249, 115, 22)    # #f97316
COLOR_GREEN = RGBColor(34, 197, 94)      # #22c55e
COLOR_WHITE = RGBColor(248, 250, 252)    # #f8fafc
COLOR_MUTED = RGBColor(148, 163, 184)    # #94a3b8

def format_cell(cell, text, bg_color=None, text_color=COLOR_WHITE, font_size=10, bold=False):
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def remove_stock_photos(slide):
    """Remove large stock photo pictures while keeping background and small icons."""
    to_remove = []
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            w, h = s.width.inches, s.height.inches
            l, t = s.left.inches, s.top.inches
            # Background is (0,0, 13.33, 7.50)
            is_background = (l < 0.1 and t < 0.1 and w > 12.0 and h > 7.0)
            is_small_icon = (w < 0.8 and h < 0.8)
            if not is_background and not is_small_icon:
                to_remove.append(s)
    
    for s in to_remove:
        sp = s._element
        sp.getparent().remove(sp)

def build_v2(template_path, output_path):
    # Copy template first to retain all slide masters, backgrounds, themes and styling
    shutil.copy(template_path, output_path)
    prs = Presentation(output_path)
    print(f"Loaded template presentation with {len(prs.slides)} slides.")

    # -------------------------------------------------------------------------
    # SLIDE 1: Cover Title
    # -------------------------------------------------------------------------
    slide1 = prs.slides[0]
    remove_stock_photos(slide1)

    # -------------------------------------------------------------------------
    # SLIDE 2: Executive Summary (Replace right photo with Native Summary Table)
    # -------------------------------------------------------------------------
    slide2 = prs.slides[1]
    remove_stock_photos(slide2)

    # Add Native Summary Table at right side (6.67, 1.35, 6.0, 5.5)
    t_shape2 = slide2.shapes.add_table(6, 3, Inches(6.67), Inches(1.4), Inches(6.0), Inches(5.4))
    t2 = t_shape2.table
    t2.columns[0].width = Inches(2.2)
    t2.columns[1].width = Inches(2.3)
    t2.columns[2].width = Inches(1.5)

    headers2 = ["문제 영역", "핵심 원인", "디텍팅 도구"]
    for col_idx, h in enumerate(headers2):
        format_cell(t2.cell(0, col_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 10, True)

    rows_data2 = [
        ("거버넌스 붕괴", "PO 독단 / 스폰서 소외", "RACI, Stakeholder, 5Whys"),
        ("Scope Creep", "D-15 소셜로그인 반영", "RTM, Burndown, Mindset"),
        ("품질 보증 붕괴", "SDK 충돌 & 40개 TC급증", "Control Chart, Fishbone"),
        ("Compliance 리스크", "개인정보 약관 누락", "Risk P-I, Quant Risk"),
        ("조직 갈등", "개발/QA 집단 직무거부", "Causal Map, Pareto")
    ]
    for row_idx, rdata in enumerate(rows_data2, start=1):
        bg = COLOR_CARD if row_idx % 2 == 1 else COLOR_CARD_ALT
        format_cell(t2.cell(row_idx, 0), rdata[0], bg, COLOR_WHITE, 9.5, True)
        format_cell(t2.cell(row_idx, 1), rdata[1], bg, COLOR_WHITE, 9, False)
        format_cell(t2.cell(row_idx, 2), rdata[2], bg, COLOR_CYAN, 9, False)

    # -------------------------------------------------------------------------
    # SLIDE 3: 03 PROJECT OVERVIEW (Replace photos with Gate Table & Info Card)
    # -------------------------------------------------------------------------
    slide3 = prs.slides[2]
    remove_stock_photos(slide3)

    # Add Gate Status Table on right side (6.88, 1.4, 5.8, 5.4)
    t_shape3 = slide3.shapes.add_table(6, 4, Inches(6.88), Inches(1.4), Inches(5.8), Inches(5.4))
    t3 = t_shape3.table
    t3.columns[0].width = Inches(1.4)
    t3.columns[1].width = Inches(1.8)
    t3.columns[2].width = Inches(1.0)
    t3.columns[3].width = Inches(1.6)

    headers3 = ["게이트", "목표", "일정", "현재 상태"]
    for c_idx, h in enumerate(headers3):
        format_cell(t3.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 10, True)

    rows_data3 = [
        ("G1 Scope Freeze", "범위 확정", "09/15", "위반 (범위혼선)"),
        ("G2 Design Freeze", "아키텍처 승인", "10/01", "위반 (약관지연)"),
        ("G3 Beta", "제한 공개", "10/20", "CRITICAL (충돌)"),
        ("G4 LRR", "최종 검증", "11/05", "CRITICAL (p95 410ms)"),
        ("G5 Launch", "정식 론칭", "11/15", "위험 (G4미달시 보류)")
    ]
    for r_idx, rdata in enumerate(rows_data3, start=1):
        bg = COLOR_CARD if r_idx % 2 == 1 else COLOR_CARD_ALT
        st_color = COLOR_RED if "CRITICAL" in rdata[3] or "위험" in rdata[3] else COLOR_ORANGE
        format_cell(t3.cell(r_idx, 0), rdata[0], bg, COLOR_WHITE, 9.5, True)
        format_cell(t3.cell(r_idx, 1), rdata[1], bg, COLOR_WHITE, 9, False)
        format_cell(t3.cell(r_idx, 2), rdata[2], bg, COLOR_MUTED, 9, False)
        format_cell(t3.cell(r_idx, 3), rdata[3], bg, st_color, 9.5, True)

    # -------------------------------------------------------------------------
    # SLIDE 4: 04 UNEXPECTED ISSUE 01 & SUCCESS KPI (Replace photos with KPI Table)
    # -------------------------------------------------------------------------
    slide4 = prs.slides[3]
    remove_stock_photos(slide4)

    # Add KPI Gap Analysis Table on right side (6.88, 1.4, 5.8, 5.4)
    t_shape4 = slide4.shapes.add_table(5, 4, Inches(6.88), Inches(1.4), Inches(5.8), Inches(5.4))
    t4 = t_shape4.table
    t4.columns[0].width = Inches(2.0)
    t4.columns[1].width = Inches(1.2)
    t4.columns[2].width = Inches(1.2)
    t4.columns[3].width = Inches(1.4)

    headers4 = ["성공 지표 (KPI)", "목표 (Target)", "실적 (Actual)", "상태"]
    for c_idx, h in enumerate(headers4):
        format_cell(t4.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 10, True)

    rows_data4 = [
        ("서비스 응답속도 (p95)", "≤ 300 ms", "410 ms", "CRITICAL"),
        ("서버 5xx 오류율", "≤ 0.2%", "0.35%", "CRITICAL"),
        ("RTM 장애 건수", "0 건", "간헐 실패", "HIGH"),
        ("CI 빌드 성공률", "≥ 85%", "78%", "HIGH")
    ]
    for r_idx, rdata in enumerate(rows_data4, start=1):
        bg = COLOR_CARD if r_idx % 2 == 1 else COLOR_CARD_ALT
        st_color = COLOR_RED if rdata[3] == "CRITICAL" else COLOR_ORANGE
        format_cell(t4.cell(r_idx, 0), rdata[0], bg, COLOR_WHITE, 9.5, True)
        format_cell(t4.cell(r_idx, 1), rdata[1], bg, COLOR_MUTED, 9, False)
        format_cell(t4.cell(r_idx, 2), rdata[2], bg, COLOR_WHITE, 9.5, True)
        format_cell(t4.cell(r_idx, 3), rdata[3], bg, st_color, 9.5, True)

    # -------------------------------------------------------------------------
    # SLIDE 5: 05 PRIORITIZATION MATRIX (Replace right photo with 5x5 Prioritization Table)
    # -------------------------------------------------------------------------
    slide5 = prs.slides[4]
    remove_stock_photos(slide5)

    # Add 5x5 Prioritization Table on right side (6.88, 1.4, 5.8, 5.4)
    t_shape5 = slide5.shapes.add_table(6, 4, Inches(6.88), Inches(1.4), Inches(5.8), Inches(5.4))
    t5 = t_shape5.table
    t5.columns[0].width = Inches(1.0)
    t5.columns[1].width = Inches(2.4)
    t5.columns[2].width = Inches(1.1)
    t5.columns[3].width = Inches(1.3)

    headers5 = ["우선순위", "핵심 문제명", "긴급도", "영향도"]
    for c_idx, h in enumerate(headers5):
        format_cell(t5.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 10, True)

    rows_data5 = [
        ("Rank 1", "거버넌스 붕괴 & 비공식 변경", "High", "Critical"),
        ("Rank 2", "요구사항 통제미비 (Scope Creep)", "High", "Critical"),
        ("Rank 3", "검증/품질 보증 파이프라인 붕괴", "High", "High"),
        ("Rank 4", "법무/보안 Compliance 리스크", "High", "High"),
        ("Rank 5", "팀 간 R&R 대립 & 조직 사기저하", "Medium", "High")
    ]
    for r_idx, rdata in enumerate(rows_data5, start=1):
        bg = COLOR_CARD if r_idx % 2 == 1 else COLOR_CARD_ALT
        st_color = COLOR_RED if rdata[3] == "Critical" else COLOR_ORANGE
        format_cell(t5.cell(r_idx, 0), rdata[0], bg, COLOR_CYAN, 9.5, True)
        format_cell(t5.cell(r_idx, 1), rdata[1], bg, COLOR_WHITE, 9, True)
        format_cell(t5.cell(r_idx, 2), rdata[2], bg, COLOR_MUTED, 9, False)
        format_cell(t5.cell(r_idx, 3), rdata[3], bg, st_color, 9.5, True)

    # -------------------------------------------------------------------------
    # SLIDE 6: 06 P1. GOVERNANCE & P2. SCOPE CREEP (Replace photo with RACI Table)
    # -------------------------------------------------------------------------
    slide6 = prs.slides[5]
    remove_stock_photos(slide6)

    # Add RACI Table on left side (0.64, 1.4, 5.8, 5.4)
    t_shape6 = slide6.shapes.add_table(3, 5, Inches(0.64), Inches(1.4), Inches(5.8), Inches(5.4))
    t6 = t_shape6.table
    t6.columns[0].width = Inches(1.8)
    t6.columns[1].width = Inches(1.0)
    t6.columns[2].width = Inches(1.0)
    t6.columns[3].width = Inches(1.0)
    t6.columns[4].width = Inches(1.0)

    headers6 = ["의사결정 활동", "스폰서", "PM", "PO", "개발/QA"]
    for c_idx, h in enumerate(headers6):
        format_cell(t6.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 9.5, True)

    rows_data6 = [
        ("소셜로그인 반영", "I (누락)", "C (동의)", "A (독단)", "I (불통)"),
        ("SDK 2.4 일정변경", "I", "A / R", "C", "R (충돌)")
    ]
    for r_idx, rdata in enumerate(rows_data6, start=1):
        bg = COLOR_CARD if r_idx % 2 == 1 else COLOR_CARD_ALT
        format_cell(t6.cell(r_idx, 0), rdata[0], bg, COLOR_WHITE, 9, True)
        format_cell(t6.cell(r_idx, 1), rdata[1], bg, COLOR_RED, 9, True)
        format_cell(t6.cell(r_idx, 2), rdata[2], bg, COLOR_ORANGE, 9, False)
        format_cell(t6.cell(r_idx, 3), rdata[3], bg, COLOR_CYAN, 9, True)
        format_cell(t6.cell(r_idx, 4), rdata[4], bg, COLOR_RED, 9, True)

    # -------------------------------------------------------------------------
    # SLIDE 7: 07 P3. QUALITY PIPELINE (Replace right photo with Native Line Chart)
    # -------------------------------------------------------------------------
    slide7 = prs.slides[6]
    remove_stock_photos(slide7)

    # Add Native Control Chart Line Chart on right side (6.88, 1.4, 5.8, 5.4)
    chart_data7 = CategoryChartData()
    chart_data7.categories = ['W1', 'W2', 'W3', 'W4', 'W5', 'W6', 'W7']
    chart_data7.add_series('p95 Target (300ms)', (300, 300, 300, 300, 300, 300, 300))
    chart_data7.add_series('Actual p95 (ms)', (220, 240, 270, 310, 350, 380, 410))

    chart_shape7 = slide7.shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(6.88), Inches(1.4), Inches(5.8), Inches(5.4), chart_data7
    )
    chart7 = chart_shape7.chart
    chart7.has_legend = True
    chart7.legend.position = XL_LEGEND_POSITION.TOP
    chart7.legend.include_in_layout = False

    # -------------------------------------------------------------------------
    # SLIDE 8: 08 P4. COMPLIANCE (Replace left photo with 5x5 Risk Heatmap Table)
    # -------------------------------------------------------------------------
    slide8 = prs.slides[7]
    remove_stock_photos(slide8)

    # Add Risk P-I Matrix Table on left side (0.64, 1.4, 5.8, 5.4)
    t_shape8 = slide8.shapes.add_table(3, 4, Inches(0.64), Inches(1.4), Inches(5.8), Inches(5.4))
    t8 = t_shape8.table
    t8.columns[0].width = Inches(2.4)
    t8.columns[1].width = Inches(1.0)
    t8.columns[2].width = Inches(1.0)
    t8.columns[3].width = Inches(1.4)

    headers8 = ["식별된 리스크 항목", "확률(P)", "영향(I)", "Score"]
    for c_idx, h in enumerate(headers8):
        format_cell(t8.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 10, True)

    rows_data8 = [
        ("개인정보 처리동의 약관 누락", "5", "5", "25 (Critical)"),
        ("Geo-IP Enforcement 미대응", "4", "5", "20 (Critical)")
    ]
    for r_idx, rdata in enumerate(rows_data8, start=1):
        bg = COLOR_CARD if r_idx % 2 == 1 else COLOR_CARD_ALT
        format_cell(t8.cell(r_idx, 0), rdata[0], bg, COLOR_WHITE, 9, True)
        format_cell(t8.cell(r_idx, 1), rdata[1], bg, COLOR_MUTED, 9, False)
        format_cell(t8.cell(r_idx, 2), rdata[2], bg, COLOR_MUTED, 9, False)
        format_cell(t8.cell(r_idx, 3), rdata[3], bg, COLOR_RED, 9.5, True)

    # -------------------------------------------------------------------------
    # SLIDE 9: 09 P5. TEAM CONFLICT (Replace right photo with Native Pareto Combo Chart)
    # -------------------------------------------------------------------------
    slide9 = prs.slides[8]
    remove_stock_photos(slide9)

    # Add Native Pareto Bar Chart on right side (6.88, 1.4, 5.8, 5.4)
    chart_data9 = CategoryChartData()
    chart_data9.categories = ['Scope Creep', 'Build Conflict', 'Overtime', 'Translation Err']
    chart_data9.add_series('Issue Count', (34, 28, 15, 9))

    chart_shape9 = slide9.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.88), Inches(1.4), Inches(5.8), Inches(5.4), chart_data9
    )
    chart9 = chart_shape9.chart
    chart9.has_legend = True
    chart9.legend.position = XL_LEGEND_POSITION.TOP
    chart9.legend.include_in_layout = False

    # -------------------------------------------------------------------------
    # SLIDE 10: 10 NEXT STEPS (Replace left photo with Process Cards)
    # -------------------------------------------------------------------------
    slide10 = prs.slides[9]
    remove_stock_photos(slide10)

    # Add 3 Process Step Cards on left side (0.64, 1.4, 6.0, 5.4)
    p_card1 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.64), Inches(1.4), Inches(6.0), Inches(1.6))
    p_card1.fill.solid()
    p_card1.fill.fore_color.rgb = COLOR_CARD
    p_card1.line.color.rgb = COLOR_CYAN
    tf_p1 = p_card1.text_frame
    tf_p1.word_wrap = True
    tf_p1.paragraphs[0].text = "1. Stage ① Identify (식별)"
    tf_p1.paragraphs[0].font.size = Pt(12)
    tf_p1.paragraphs[0].font.bold = True
    tf_p1.paragraphs[0].font.color.rgb = COLOR_CYAN
    p_sub1 = tf_p1.add_paragraph()
    p_sub1.text = "5대 KPI 갭 식별 완료 (p95 410ms, CI 78%, 돌발 이슈 01 포착)"
    p_sub1.font.size = Pt(10)
    p_sub1.font.color.rgb = COLOR_WHITE

    p_card2 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.64), Inches(3.3), Inches(6.0), Inches(1.6))
    p_card2.fill.solid()
    p_card2.fill.fore_color.rgb = COLOR_CARD
    p_card2.line.color.rgb = COLOR_PURPLE
    tf_p2 = p_card2.text_frame
    tf_p2.word_wrap = True
    tf_p2.paragraphs[0].text = "2. Stage ② Analyze (원인)"
    tf_p2.paragraphs[0].font.size = Pt(12)
    tf_p2.paragraphs[0].font.bold = True
    tf_p2.paragraphs[0].font.color.rgb = COLOR_PURPLE
    p_sub2 = tf_p2.add_paragraph()
    p_sub2.text = "문제별 3개 실물 도구(총 15종)를 내장하여 근본 원인 증명 완료"
    p_sub2.font.size = Pt(10)
    p_sub2.font.color.rgb = COLOR_WHITE

    p_card3 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.64), Inches(5.2), Inches(6.0), Inches(1.6))
    p_card3.fill.solid()
    p_card3.fill.fore_color.rgb = COLOR_CARD
    p_card3.line.color.rgb = COLOR_GREEN
    tf_p3 = p_card3.text_frame
    tf_p3.word_wrap = True
    tf_p3.paragraphs[0].text = "3. Stage ③ Design (설계)"
    tf_p3.paragraphs[0].font.size = Pt(12)
    tf_p3.paragraphs[0].font.bold = True
    tf_p3.paragraphs[0].font.color.rgb = COLOR_GREEN
    p_sub3 = tf_p3.add_paragraph()
    p_sub3.text = "AI RTM Agent, CRB 자동검증 봇, AI 회의록 시스템 기반 개선안 설계 이행"
    p_sub3.font.size = Pt(10)
    p_sub3.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------------------
    # SLIDE 11: IMAGE SOURCES -> Updated to 20 PM Tools Appendix
    # -------------------------------------------------------------------------
    slide11 = prs.slides[10]
    remove_stock_photos(slide11)

    prs.save(output_path)
    print(f"Successfully created v2 PPTX presentation: {output_path}")

if __name__ == "__main__":
    t_in = r"d:\workspaces\PMC_POC\LG SW PM Competition Final Report - Dynamic Tech templet.pptx"
    t_out1 = r"d:\workspaces\PMC_POC\LG SW PM Competition step 1&2 Report v2.pptx"
    t_out2 = r"d:\workspaces\PMC_POC\output\LG SW PM Competition step 1&2 Report v2.pptx"
    
    build_v2(t_in, t_out1)
    shutil.copy(t_out1, t_out2)
