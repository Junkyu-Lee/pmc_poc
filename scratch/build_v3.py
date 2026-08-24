import sys
import os
import shutil
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

sys.stdout.reconfigure(encoding='utf-8')

# Colors
COLOR_BG_DARK = RGBColor(15, 23, 42)
COLOR_CARD = RGBColor(30, 41, 59)
COLOR_CARD_ALT = RGBColor(51, 65, 85)
COLOR_CYAN = RGBColor(56, 189, 248)
COLOR_PURPLE = RGBColor(168, 85, 247)
COLOR_RED = RGBColor(244, 63, 94)
COLOR_ORANGE = RGBColor(249, 115, 22)
COLOR_GREEN = RGBColor(34, 197, 94)
COLOR_WHITE = RGBColor(248, 250, 252)
COLOR_MUTED = RGBColor(148, 163, 184)

def format_cell(cell, text, bg_color=None, text_color=COLOR_WHITE, font_size=10, bold=False, align=PP_ALIGN.CENTER):
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_text_box(slide, text, left, top, width, height, font_size=12, color=COLOR_WHITE, bold=False, wrap=True, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def duplicate_slide(prs, source_index):
    source_slide = prs.slides[source_index]
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_el)
    for rel in source_slide.part.rels.values():
        if 'image' in rel.reltype:
            new_slide.part.relate_to(rel.target_part, rel.reltype)
    return new_slide

def remove_stock_photos(slide):
    to_remove = []
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            w, h = s.width.inches, s.height.inches
            l, t = s.left.inches, s.top.inches
            is_background = (l < 0.1 and t < 0.1 and w > 12.0 and h > 7.0)
            is_small_icon = (w < 1.5 and h < 1.5)
            if not is_background and not is_small_icon:
                to_remove.append(s)
    for s in to_remove:
        sp = s._element
        sp.getparent().remove(sp)

def build_v3(tpl_path, output_path):
    prs = Presentation(tpl_path)
    
    # Target 19 slides:
    # We will duplicate from existing 11 slides
    # Orig 0..5 -> New 0..5
    # Orig 6 (P1) -> New 6, 7 (P1 A, P1 B)
    #             -> New 8, 9 (P2 A, P2 B)
    #             -> New 10, 11 (P3 A, P3 B)
    # Orig 7 (P4) -> New 12, 13 (P4 A, P4 B)
    # Orig 8 (P5) -> New 14, 15 (P5 A, P5 B)
    # Orig 9 (Next)-> New 16
    # Orig 10 (App)-> New 17
    # Orig 0 (Cov)-> New 18
    
    # To avoid index shifting mess, we will copy the entire presentation to a new list of slides
    # Actually, we can just duplicate to the end, then move.
    # But since pptx doesn't support move, we will create a NEW presentation and copy slides into it.
    # Wait, `prs.slides.add_slide` adds to the end.
    
    prs_out = Presentation(tpl_path)
    
    # We have 11 slides initially.
    # We want to create slides at the end.
    # 11: P1 B (dup 6)
    s11 = duplicate_slide(prs_out, 6)
    # 12: P2 A (dup 6)
    s12 = duplicate_slide(prs_out, 6)
    # 13: P2 B (dup 6)
    s13 = duplicate_slide(prs_out, 6)
    # 14: P3 A (dup 6)
    s14 = duplicate_slide(prs_out, 6)
    # 15: P3 B (dup 6)
    s15 = duplicate_slide(prs_out, 6)
    # 16: P4 B (dup 7)
    s16 = duplicate_slide(prs_out, 7)
    # 17: P5 B (dup 8)
    s17 = duplicate_slide(prs_out, 8)
    # 18: Thank You (dup 0)
    s18 = duplicate_slide(prs_out, 0)
    
    # Now map the logical 19 slides to the physical slides in prs_out:
    s = [
        prs_out.slides[0],  # 0 Cover
        prs_out.slides[1],  # 1 Exec Summary
        prs_out.slides[2],  # 2 Proj Overview
        prs_out.slides[3],  # 3 Unexpected Issue
        prs_out.slides[4],  # 4 KPI Gap
        prs_out.slides[5],  # 5 Prioritization
        prs_out.slides[6],  # 6 P1 A
        prs_out.slides[11], # 7 P1 B
        prs_out.slides[12], # 8 P2 A
        prs_out.slides[13], # 9 P2 B
        prs_out.slides[14], # 10 P3 A
        prs_out.slides[15], # 11 P3 B
        prs_out.slides[7],  # 12 P4 A
        prs_out.slides[16], # 13 P4 B
        prs_out.slides[8],  # 14 P5 A
        prs_out.slides[17], # 15 P5 B
        prs_out.slides[9],  # 16 Bridge
        prs_out.slides[10], # 17 Appendix
        prs_out.slides[18]  # 18 Thank You
    ]

    # --- Slide 0: Cover ---
    remove_stock_photos(s[0])
    # Title is already good, subtitle we can adjust
    add_text_box(s[0], "시나리오: NovaHome Connect 통합 프로젝트 + 돌발 이슈 01 (내부 갈등: 베타 D-15 소셜 로그인 강행 및 대립)", Inches(1), Inches(4.5), Inches(8), Inches(1), 14, COLOR_CYAN, True)

    # --- Slide 1: Exec Summary ---
    remove_stock_photos(s[1])
    add_text_box(s[1], "핵심 진단: 현상은 '일정 지연 및 직무 대립'이나, 근본 원인은 PMBOK 8판 체계 미비, 변경 통제(CRB) 부재, 그리고 비공식적 의사결정 방식에 있습니다.", Inches(6.8), Inches(1.5), Inches(5.5), Inches(1.2), 11, COLOR_RED, True)
    
    t_shape1 = s[1].shapes.add_table(6, 3, Inches(6.8), Inches(2.8), Inches(5.5), Inches(3.5))
    t1 = t_shape1.table
    t1.columns[0].width = Inches(1.5)
    t1.columns[1].width = Inches(2.2)
    t1.columns[2].width = Inches(1.8)
    for c_idx, h in enumerate(["문제 영역", "핵심 내용", "디텍팅 도구 (3종)"]):
        format_cell(t1.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 9, True)
    data1 = [
        ("거버넌스 붕괴", "PO 독단 / 스폰서 소외", "RACI, Stakeholder, 5Whys"),
        ("Scope Creep", "D-15 소셜로그인 반영", "RTM, Burndown, Mindset"),
        ("품질 보증 붕괴", "SDK 충돌 & 40TC급증", "Control Chart, Fishbone"),
        ("Compliance", "개인정보 약관 누락", "Risk P-I, Quant Risk"),
        ("조직 갈등", "개발/QA 집단 직무거부", "Causal Map, Pareto")
    ]
    for r, row in enumerate(data1, 1):
        for c, val in enumerate(row):
            format_cell(t1.cell(r, c), val, COLOR_CARD if r%2==1 else COLOR_CARD_ALT, COLOR_WHITE, 8)

    # --- Slide 2: Project Overview ---
    remove_stock_photos(s[2])
    t_shape2 = s[2].shapes.add_table(6, 4, Inches(6.88), Inches(1.4), Inches(5.8), Inches(5.4))
    t2 = t_shape2.table
    for c_idx, h in enumerate(["게이트", "목표", "일정", "상태 & 리스크"]):
        format_cell(t2.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 9, True)
    data2 = [
        ("G1 Scope Freeze", "범위 확정", "09/15", "위반 (범위혼선)"),
        ("G2 Design Freeze", "아키텍처 승인", "10/01", "위반 (약관지연)"),
        ("G3 Beta", "제한공개", "10/20", "CRITICAL (소셜로그인 충돌)"),
        ("G4 LRR", "최종검증", "11/05", "CRITICAL (p95 410ms)"),
        ("G5 Launch", "정식론칭", "11/15", "위험 (미해소시 실패)")
    ]
    for r, row in enumerate(data2, 1):
        for c, val in enumerate(row):
            tc = COLOR_RED if "CRITICAL" in val or "위험" in val else (COLOR_ORANGE if "위반" in val else COLOR_WHITE)
            format_cell(t2.cell(r, c), val, COLOR_CARD if r%2==1 else COLOR_CARD_ALT, tc, 8, bold=(c==3))

    # --- Slide 3: Issue 01 ---
    remove_stock_photos(s[3])
    add_text_box(s[3], "PO / 마케팅: '마케팅 캠페인 연동 필수. 통합 테스트 강행해야 함.'\n\n서비스 BE (오지후): 'SDK 2.4 빌드 충돌 발생. 자원/QA 지원 불가.'\n\n품질/QA (문해인): '테스트 매트릭스 2배 급증. 전수 검증 불가능.'\n\n보안/법무 (윤세라): '약관 초안 누락. CRB 승인 불가시 론칭 불가.'", Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.5), 11, COLOR_WHITE, False)

    # --- Slide 4: KPI Gap Analysis ---
    remove_stock_photos(s[4])
    t_shape4 = s[4].shapes.add_table(5, 4, Inches(6.88), Inches(1.4), Inches(5.8), Inches(2.5))
    t4 = t_shape4.table
    for c_idx, h in enumerate(["KPI 지표", "목표", "실적", "상태"]):
        format_cell(t4.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 9, True)
    data4 = [
        ("응답속도 (p95)", "≤ 300 ms", "410 ms", "CRITICAL"),
        ("서버 5xx 오류율", "≤ 0.2%", "0.35%", "CRITICAL"),
        ("RTM 장애 건수", "0 건", "간헐적 실패", "HIGH"),
        ("CI 빌드 성공률", "≥ 85%", "78%", "HIGH")
    ]
    for r, row in enumerate(data4, 1):
        for c, val in enumerate(row):
            tc = COLOR_RED if val=="CRITICAL" else (COLOR_ORANGE if val=="HIGH" else COLOR_WHITE)
            format_cell(t4.cell(r, c), val, COLOR_CARD if r%2==1 else COLOR_CARD_ALT, tc, 8, bold=(c==3))

    # --- Slide 5: Prioritization ---
    remove_stock_photos(s[5])
    t_shape5 = s[5].shapes.add_table(6, 4, Inches(6.88), Inches(1.4), Inches(5.8), Inches(4.0))
    t5 = t_shape5.table
    for c_idx, h in enumerate(["우선순위", "핵심 문제명", "긴급도", "영향도"]):
        format_cell(t5.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 9, True)
    data5 = [
        ("Rank 1", "거버넌스 붕괴 & 비공식 변경", "High", "Critical"),
        ("Rank 2", "요구사항 통제미비 (Scope Creep)", "High", "Critical"),
        ("Rank 3", "검증/품질 보증 붕괴", "High", "High"),
        ("Rank 4", "법무/보안 Compliance 리스크", "High", "High"),
        ("Rank 5", "팀 간 R&R 대립 & 사기저하", "Medium", "High")
    ]
    for r, row in enumerate(data5, 1):
        for c, val in enumerate(row):
            tc = COLOR_RED if val=="Critical" else (COLOR_ORANGE if val=="High" else COLOR_WHITE)
            format_cell(t5.cell(r, c), val, COLOR_CARD if r%2==1 else COLOR_CARD_ALT, tc, 8, bold=(c==3))

    # Helper for Problem A slides (Analysis)
    def make_problem_a(slide, title, mindset, domain, aicase, future):
        remove_stock_photos(slide)
        add_text_box(slide, title, Inches(0.5), Inches(0.5), Inches(10), Inches(0.8), 20, COLOR_CYAN, True)
        content = f"1. Mindset: {mindset}\n\n2. Domain: {domain}\n\n3. AI Use Case: {aicase}\n\n4. Future State: {future}"
        add_text_box(slide, content, Inches(1.0), Inches(2.0), Inches(11), Inches(4.5), 14, COLOR_WHITE, False)

    # --- Slide 6: P1 A ---
    make_problem_a(s[6], "🔍 [문제 1] 거버넌스 붕괴 & 권한/스폰서십 비공식 변경 (분석)",
                   "의사결정을 공식 체계(CRB)가 아닌 '영업/마케팅 단독 압박 및 임기응변식 타협'으로 처리함.",
                   "Governance Performance Domain & Stakeholder Performance Domain 붕괴.",
                   "AI 기반 의사결정 영향도 자동 검증 & RACI 승인 에이전트 구축.",
                   "스폰서-PO-PM 3자 공식 승인 절차 미거친 변경 요청의 투입 0건 달성.")
    
    # --- Slide 7: P1 B (Tools) ---
    remove_stock_photos(s[7])
    add_text_box(s[7], "📊 [문제 1] 디텍팅 실물 도구 3종 풀패키지", Inches(0.5), Inches(0.5), Inches(10), Inches(0.8), 20, COLOR_CYAN, True)
    # RACI Table
    add_text_box(s[7], "Tool #3: RACI Matrix", Inches(0.5), Inches(1.5), Inches(3), Inches(0.5), 12, COLOR_CYAN, True)
    t_shape7_1 = s[7].shapes.add_table(3, 5, Inches(0.5), Inches(2.0), Inches(5.5), Inches(1.5))
    t7_1 = t_shape7_1.table
    for c_idx, h in enumerate(["활동", "스폰서", "PM", "PO", "개발/QA"]): format_cell(t7_1.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 8, True)
    for r, row in enumerate([("소셜로그인", "I", "C", "A (독단)", "I"), ("SDK 2.4", "I", "A/R", "C", "R (충돌)")], 1):
        for c, val in enumerate(row): format_cell(t7_1.cell(r, c), val, COLOR_CARD if r%2==1 else COLOR_CARD_ALT, COLOR_WHITE, 8)
    
    # Stakeholder Table
    add_text_box(s[7], "Tool #13: Stakeholder Engagement", Inches(7.0), Inches(1.5), Inches(4), Inches(0.5), 12, COLOR_CYAN, True)
    t_shape7_2 = s[7].shapes.add_table(4, 4, Inches(7.0), Inches(2.0), Inches(5.5), Inches(2.0))
    t7_2 = t_shape7_2.table
    for c_idx, h in enumerate(["이해관계자", "현재(C)", "바람직(D)", "갭"]): format_cell(t7_2.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 8, True)
    for r, row in enumerate([("스폰서", "Unaware", "Leading", "단절"), ("마케팅", "Resistant", "Supportive", "대립"), ("QA", "Resistant", "Supportive", "거부")], 1):
        for c, val in enumerate(row): format_cell(t7_2.cell(r, c), val, COLOR_CARD if r%2==1 else COLOR_CARD_ALT, COLOR_WHITE, 8)
    
    add_text_box(s[7], "Tool #9: 5 Whys Root Cause Flow\n\n1. 기습 반영 ➔ 2. PO 일방수용 ➔ 3. PM 미제지 ➔ 4. 통제불가 ➔ 5. Root Cause: PMBOK 거버넌스 부재", Inches(0.5), Inches(4.5), Inches(12), Inches(2), 12, COLOR_ORANGE, True)

    # --- Slide 8: P2 A ---
    make_problem_a(s[8], "🔍 [문제 2] 요구사항 변경 통제 미비 & Scope Creep (분석)",
                   "G1 Scope Freeze 베이스라인을 경시하고 기술적 의존관계를 통제하지 않음.",
                   "Planning Performance Domain & Scope Management Domain 붕괴.",
                   "AI RTM 추적성 분석기 & 스코프 변동 자동 경고 봇 구축.",
                   "베이스라인 변경 시 자동 영향도 분석 및 스코프 크립 0건 달성.")

    # --- Slide 9: P2 B ---
    remove_stock_photos(s[9])
    add_text_box(s[9], "📊 [문제 2] 디텍팅 실물 도구 3종 풀패키지", Inches(0.5), Inches(0.5), Inches(10), Inches(0.8), 20, COLOR_CYAN, True)
    add_text_box(s[9], "Tool #2: RTM Traceability\nREQ-017(소셜로그인) ➔ REQ-021약관 의존성 파괴\nREQ-011(RTM알림) ➔ SDK 2.4 Blocked", Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.0), 12, COLOR_WHITE, False)
    
    add_text_box(s[9], "Tool #18: Sprint 3 Burndown", Inches(7.0), Inches(1.5), Inches(4), Inches(0.5), 12, COLOR_CYAN, True)
    cd9 = CategoryChartData()
    cd9.categories = ['D1', 'D2', 'D3', 'D4', 'D5', 'D6', 'D7']
    cd9.add_series('Target', (100, 85, 70, 55, 40, 25, 10))
    cd9.add_series('Actual (Scope Creep)', (100, 85, 75, 75, 70, 85, 70))
    s[9].shapes.add_chart(XL_CHART_TYPE.LINE, Inches(7.0), Inches(2.0), Inches(5.5), Inches(2.5), cd9)
    
    add_text_box(s[9], "Tool #11: Mindset Mapping\n기존: '요구사항 추가는 흔한 일, 개발자가 버티면 됨'\nPMBOK: 'Scope Freeze 준수 및 영향도 자동 통제 체계 수립'", Inches(0.5), Inches(4.5), Inches(12), Inches(1.5), 12, COLOR_ORANGE, True)

    # --- Slide 10: P3 A ---
    make_problem_a(s[10], "🔍 [문제 3] 검증/품질 보증 파이프라인 붕괴 (분석)",
                   "단기 임시 대책이 전체 테스트 파이프라인을 붕괴시키는 품질 부채 간과.",
                   "Delivery Performance Domain & Quality Performance Domain 붕괴.",
                   "AI 자동 테스트 케이스 생성기 & 회귀 성능 예측 모니터링 에이전트.",
                   "CI 빌드 성공률 90% 이상 확보 및 p95 ≤ 300ms 목표 안정 달성.")

    # --- Slide 11: P3 B ---
    remove_stock_photos(s[11])
    add_text_box(s[11], "📊 [문제 3] 디텍팅 실물 도구 3종 풀패키지", Inches(0.5), Inches(0.5), Inches(10), Inches(0.8), 20, COLOR_CYAN, True)
    
    add_text_box(s[11], "Tool #17: Control Chart (SVG/Line)", Inches(0.5), Inches(1.5), Inches(4), Inches(0.5), 12, COLOR_CYAN, True)
    cd11 = CategoryChartData()
    cd11.categories = ['W1', 'W2', 'W3', 'W4', 'W5', 'W6', 'W7']
    cd11.add_series('p95 LCL', (300, 300, 300, 300, 300, 300, 300))
    cd11.add_series('Actual p95', (220, 240, 270, 310, 350, 380, 410))
    s[11].shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.5), Inches(2.0), Inches(5.5), Inches(2.5), cd11)
    
    add_text_box(s[11], "Tool #8: Variance Analysis\n응답속도 p95: 410ms (+110ms 분산)\nCI 성공률: 78% (-7%p 분산)", Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.0), 12, COLOR_WHITE, False)
    add_text_box(s[11], "Tool #10: Fishbone Diagram\nPeople: QA 40TC 전수불능 | Process: 캐시 임시적용 | Tech: SDK 충돌 ➔ CI 78% 추락", Inches(0.5), Inches(5.0), Inches(12), Inches(1.5), 12, COLOR_ORANGE, True)

    # --- Slide 12: P4 A ---
    make_problem_a(s[12], "🔍 [문제 4] 법무/보안 규제 준수(Compliance) 리스크 (분석)",
                   "보안/법무 약관을 프로젝트 론칭 직전의 형식적 절차로 치부하는 민감도 결여.",
                   "Risk Performance Domain & Compliance Domain 붕괴.",
                   "AI 컴플라이언스 사전 약관 검증기 & Geo-IP 자동 신청 에이전트.",
                   "CRB 승인 100% 사전 확보 및 글로벌 법적 과징금 리스크 ZERO.")

    # --- Slide 13: P4 B ---
    remove_stock_photos(s[13])
    add_text_box(s[13], "📊 [문제 4] 디텍팅 실물 도구 3종 풀패키지", Inches(0.5), Inches(0.5), Inches(10), Inches(0.8), 20, COLOR_CYAN, True)
    add_text_box(s[13], "Tool #4: Risk P-I Heatmap", Inches(0.5), Inches(1.5), Inches(4), Inches(0.5), 12, COLOR_CYAN, True)
    t_shape13 = s[13].shapes.add_table(3, 4, Inches(0.5), Inches(2.0), Inches(5.5), Inches(1.5))
    t13 = t_shape13.table
    for c_idx, h in enumerate(["리스크 항목", "P", "I", "Score"]): format_cell(t13.cell(0, c_idx), h, COLOR_CARD_ALT, COLOR_CYAN, 8, True)
    for r, row in enumerate([("개인정보 약관 누락", "5", "5", "25 (Crit)"), ("Geo-IP 미대응", "4", "5", "20 (Crit)")], 1):
        for c, val in enumerate(row): format_cell(t13.cell(r, c), val, COLOR_CARD if r%2==1 else COLOR_CARD_ALT, COLOR_WHITE, 8)
        
    add_text_box(s[13], "Tool #16: Quant Risk Impact\n- 유럽 지역 차단\n- 과징금 정량: 매출대비 패널티\n- 론칭 불가: CRB 미통과", Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.0), 12, COLOR_WHITE, False)
    add_text_box(s[13], "Tool #14: Feedback Loop\n약관 검토 지연 ➔ CRB 승인 실패 ➔ 론칭 연기 ➔ 마케팅 손실 (악순환 반복)", Inches(0.5), Inches(4.5), Inches(12), Inches(1.5), 12, COLOR_ORANGE, True)

    # --- Slide 14: P5 A ---
    make_problem_a(s[14], "🔍 [문제 5] 팀 간 R&R 대립 및 조직 사기 저하 (분석)",
                   "리더십을 일방 통제로 인식하고 팀원과의 투명한 소통 제공 실패.",
                   "Team Performance Domain & Culture Domain 붕괴.",
                   "AI 다국어 회의록/의사결정 요약 봇 & 팀 사기 감지 에이전트.",
                   "의사결정의 100% 투명 공유 및 글로벌 팀 협업 만족도 90% 달성.")

    # --- Slide 15: P5 B ---
    remove_stock_photos(s[15])
    add_text_box(s[15], "📊 [문제 5] 디텍팅 실물 도구 3종 풀패키지", Inches(0.5), Inches(0.5), Inches(10), Inches(0.8), 20, COLOR_CYAN, True)
    add_text_box(s[15], "Tool #15: Causal Relationship Map\n거버넌스붕괴 ➔ Scope Creep ➔ 품질/약관붕괴 ➔ 팀대립/퇴사위험", Inches(0.5), Inches(1.5), Inches(5.5), Inches(1.5), 12, COLOR_WHITE, False)
    
    add_text_box(s[15], "Tool #19: Pareto Chart", Inches(7.0), Inches(1.5), Inches(4), Inches(0.5), 12, COLOR_CYAN, True)
    cd15 = CategoryChartData()
    cd15.categories = ['Scope Creep', 'Build Conflict', 'Overtime', 'Translation Err']
    cd15.add_series('Count', (34, 28, 15, 9))
    s[15].shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(2.0), Inches(5.5), Inches(2.5), cd15)
    
    add_text_box(s[15], "Tool #20: Trend Analysis\n해외 QA 번역오해 발생률 증가, 블라인드 불만 폭발 ➔ 동력 상실 위험", Inches(0.5), Inches(4.5), Inches(12), Inches(1.5), 12, COLOR_ORANGE, True)

    # --- Slide 16: Bridge ---
    remove_stock_photos(s[16])
    add_text_box(s[16], "🚀 Stage ① & ② 종합 결론 및 Stage ③ Design 이행 로드맵\n\n종합 진단 총평: 5개 문제에 매핑된 실물 분석 도구 15종을 통해 근본 원인을 증명했습니다. 이제 Value-Driven AI Agentic 시스템 설계로 즉시 이행합니다.", Inches(0.5), Inches(1.5), Inches(12), Inches(2.5), 14, COLOR_WHITE, True)

    # --- Slide 17: Appendix ---
    remove_stock_photos(s[17])
    add_text_box(s[17], "📚 Appendix 1 — 사용된 20개 PM 분석 도구 세부 목록\n\n1. KPI Gap Analysis\n2. RTM\n3. RACI\n4. Risk Matrix\n5. Symptom Mapping\n6. Impact-Urgency Matrix\n7. EVM\n8. Variance Analysis\n9. 5 Whys\n10. Fishbone\n11. Mindset Mapping\n12. Performance Domain Mapping\n13. Stakeholder Matrix\n14. Feedback Loop\n15. Causal Map\n16. Quant Risk\n17. Control Chart\n18. Burndown\n19. Pareto\n20. Trend Analysis", Inches(0.5), Inches(1.5), Inches(12), Inches(5.5), 11, COLOR_WHITE, False)

    # --- Slide 18: Thank You ---
    remove_stock_photos(s[18])
    add_text_box(s[18], "감사합니다 (Thank You)\nLG SW PM Competition 2025 Final\nNext Step: Stage ③ Design & Stage ④ Develop POC", Inches(1.0), Inches(3.5), Inches(11), Inches(2.0), 16, COLOR_CYAN, True, align=PP_ALIGN.CENTER)

    prs_out.save(output_path)
    print(f"Successfully saved {output_path} with 19 slides based on v4 HTML!")

if __name__ == "__main__":
    tpl = r"d:\workspaces\PMC_POC\LG SW PM Competition Final Report - Dynamic Tech Template.pptx"
    out1 = r"d:\workspaces\PMC_POC\LG SW PM Competition step 1&2 Report v3.pptx"
    out2 = r"d:\workspaces\PMC_POC\output\LG SW PM Competition step 1&2 Report v3.pptx"
    build_v3(tpl, out1)
    shutil.copy(out1, out2)
