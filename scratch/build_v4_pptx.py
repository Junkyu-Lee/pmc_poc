import sys
import os
import copy
import asyncio
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from playwright.async_api import async_playwright

sys.stdout.reconfigure(encoding='utf-8')

def duplicate_slide(prs, source_index):
    source_slide = prs.slides[source_index]
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_el)
    for rel in source_slide.part.rels.values():
        if 'image' in rel.reltype:
            new_slide.part.relate_to(rel.target_part, rel.reltype)
    return new_slide

def clear_all_shapes_except_background(slide):
    # We want to remove all text boxes, charts, and large background images if they are placeholders.
    # To be safe, we just remove everything except the slide layout/master elements.
    # Actually, in duplicate_slide, we copied ALL shapes. Let's remove them to have a clean slate.
    # But wait, if we remove all shapes, we might remove decorative template borders.
    # Let's only remove text shapes and placeholders.
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame or shape.is_placeholder:
            shapes_to_remove.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Remove large pictures (likely stock photos)
            w, h = shape.width.inches, shape.height.inches
            if w > 2.0 and h > 2.0:
                shapes_to_remove.append(shape)
                
    for shape in shapes_to_remove:
        try:
            element = shape.element
            element.getparent().remove(element)
        except Exception:
            pass

async def capture_html_slides():
    html_path = r"file:///d:/workspaces/PMC_POC/stage1_stage2_analysis_v4.html"
    scratch_dir = r"d:\workspaces\PMC_POC\scratch"
    screenshots = []
    
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page(
            viewport={'width': 1280, 'height': 720}, 
            device_scale_factor=2
        )
        await page.goto(html_path, wait_until='networkidle')
        
        # Inject CSS to make background transparent and hide scrollbars
        await page.add_style_tag(content="""
            body { background: transparent !important; } 
            .container { background: transparent !important; box-shadow: none !important; border: none !important; padding: 0 !important; max-width: none !important;}
            /* Let's make slide-card slightly transparent to blend with PPTX background */
            .slide-card { background: rgba(30, 41, 59, 0.95) !important; border-radius: 12px !important; margin-bottom: 2rem !important; }
        """)
        
        # Wait a bit for fonts and Base64 images to render completely
        await page.wait_for_timeout(2000)
        
        cards = await page.query_selector_all('.slide-card')
        
        for i, card in enumerate(cards):
            filename = os.path.join(scratch_dir, f'v4_slide_{i:02d}.png')
            await card.screenshot(path=filename, omit_background=True)
            screenshots.append(filename)
            
        await browser.close()
        
    return screenshots

def build_pptx(screenshots):
    template_path = r"d:\workspaces\PMC_POC\LG SW PM Competition Final Report - Dynamic Tech Template.pptx"
    output_path = r"d:\workspaces\PMC_POC\output\LG SW PM Competition step 1&2 Report v4.pptx"
    
    prs = Presentation(template_path)
    
    # We will use slide[1] (index 1) as the master template for content
    # Duplicate it for each screenshot
    content_source_idx = 1
    
    for idx, img_path in enumerate(screenshots):
        new_slide = duplicate_slide(prs, content_source_idx)
        clear_all_shapes_except_background(new_slide)
        
        # Calculate aspect ratio and center the image
        # Standard widescreen is 13.33 x 7.5 inches
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        
        # We leave a small margin
        margin = Inches(0.4)
        max_w = slide_w - (2 * margin)
        max_h = slide_h - (2 * margin)
        
        # We will just add the picture and let PPTX scale it proportionally to fit within width/height
        # Wait, python-pptx doesn't auto-scale both to fit. We just specify width, and if it exceeds height, we specify height.
        # But we can just specify left, top, width and let it keep aspect ratio.
        # A simple way to center it:
        pic = new_slide.shapes.add_picture(img_path, 0, 0)
        
        # Calculate scaling factor
        scale_w = max_w / pic.width
        scale_h = max_h / pic.height
        scale = min(scale_w, scale_h)
        
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        
        pic.left = int((slide_w - pic.width) / 2)
        pic.top = int((slide_h - pic.height) / 2)

    # Now remove the original template slides (except Title slide at index 0 if we want to keep it? No, let's keep it)
    # The original template might have 3 slides. We want to delete index 1, 2, ... up to the start of our new slides.
    # To delete a slide in python-pptx, we must remove it from the XML tree.
    xml_slides = prs.slides._sldIdLst
    slides_to_remove = []
    # Count how many original slides there were
    original_count = len(prs.slides) - len(screenshots)
    for i in range(1, original_count):  # Keep index 0 (Title)
        slides_to_remove.append(xml_slides[i])
        
    for sldId in slides_to_remove:
        xml_slides.remove(sldId)
        
    prs.save(output_path)
    print(f"Successfully generated {output_path} with {len(screenshots)} content slides!")

def main():
    print("Capturing HTML slides...")
    screenshots = asyncio.run(capture_html_slides())
    print(f"Captured {len(screenshots)} screenshots.")
    
    print("Building PPTX...")
    build_pptx(screenshots)
    
if __name__ == '__main__':
    main()
