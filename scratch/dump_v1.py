import sys, os
from pptx import Presentation

sys.stdout.reconfigure(encoding='utf-8')

def dump_v1(filepath):
    prs = Presentation(filepath)
    print(f"=== V1 FILE DUMP ({len(prs.slides)} Slides) ===")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                print(f"  [{shape.name}]: {shape.text.strip()}")

if __name__ == "__main__":
    dump_v1(r"d:\workspaces\PMC_POC\LG SW PM Competition step 1&2 Report v1.pptx")
