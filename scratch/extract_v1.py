import json
from pptx import Presentation

def extract_v1():
    prs = Presentation(r'd:\workspaces\PMC_POC\LG SW PM Competition step 1&2 Report v1.pptx')
    slides_data = []
    
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_text.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    slide_text.append(" | ".join(row_text))
        slides_data.append(slide_text)
    
    with open(r'd:\workspaces\PMC_POC\scratch\v1_content.json', 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, ensure_ascii=False, indent=2)

if __name__ == '__main__':
    extract_v1()
