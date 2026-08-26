import os, sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.stdout.reconfigure(encoding='utf-8')

def inspect_file(filepath):
    print(f"\n==================== FILE: {os.path.basename(filepath)} ====================")
    prs = Presentation(filepath)
    print(f"Slide Width: {prs.slide_width.inches} in, Height: {prs.slide_height.inches} in, Total Slides: {len(prs.slides)}")
    
    for idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {idx+1} ---")
        for s_idx, shape in enumerate(slide.shapes):
            stype = shape.shape_type
            name = shape.name
            text = shape.text.strip().replace('\n', ' ') if shape.has_text_frame and shape.text else ""
            if len(text) > 80:
                text = text[:80] + "..."
            
            p_info = ""
            if stype == MSO_SHAPE_TYPE.PICTURE:
                p_info = " [PICTURE]"
            elif stype == MSO_SHAPE_TYPE.TABLE:
                p_info = f" [TABLE {len(shape.table.rows)}x{len(shape.table.columns)}]"
            
            left, top, width, height = shape.left.inches, shape.top.inches, shape.width.inches, shape.height.inches
            print(f"  Shape {s_idx+1}: name='{name}', type={stype}{p_info}, pos=({left:.2f}, {top:.2f}, {width:.2f}, {height:.2f}), text='{text}'")

if __name__ == "__main__":
    t1 = r"d:\workspaces\PMC_POC\LG SW PM Competition Final Report - Dynamic Tech templet.pptx"
    t2 = r"d:\workspaces\PMC_POC\LG SW PM Competition step 1&2 Report v1.pptx"
    inspect_file(t1)
    inspect_file(t2)
