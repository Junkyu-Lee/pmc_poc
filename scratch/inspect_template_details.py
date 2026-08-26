import os, sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.stdout.reconfigure(encoding='utf-8')

def analyze_template(filepath):
    prs = Presentation(filepath)
    print(f"=== TEMPLATE: {os.path.basename(filepath)} ===")
    print(f"Dimensions: {prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} inches, Total Slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n==================== SLIDE {i+1} ====================")
        for s in slide.shapes:
            left, top, w, h = s.left.inches, s.top.inches, s.width.inches, s.height.inches
            stype = s.shape_type
            name = s.name
            
            if stype == MSO_SHAPE_TYPE.PICTURE:
                print(f"  📷 [PICTURE] '{name}' pos=({left:.2f}, {top:.2f}, w={w:.2f}, h={h:.2f})")
            elif stype == MSO_SHAPE_TYPE.TABLE:
                print(f"  📊 [TABLE] '{name}' pos=({left:.2f}, {top:.2f}, w={w:.2f}, h={h:.2f}) rows={len(s.table.rows)}, cols={len(s.table.columns)}")
            else:
                txt = s.text.strip().replace('\n', ' ') if s.has_text_frame and s.text else ""
                if len(txt) > 70:
                    txt = txt[:70] + "..."
                if txt:
                    print(f"  📝 [TEXT/SHAPE] '{name}' type={stype} pos=({left:.2f}, {top:.2f}, w={w:.2f}, h={h:.2f}) text='{txt}'")
                else:
                    print(f"  📐 [SHAPE] '{name}' type={stype} pos=({left:.2f}, {top:.2f}, w={w:.2f}, h={h:.2f})")

if __name__ == "__main__":
    analyze_template(r"d:\workspaces\PMC_POC\LG SW PM Competition Final Report - Dynamic Tech templet.pptx")
